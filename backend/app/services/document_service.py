"""
Document Processing Service

This service handles document upload, text extraction, chunking,
and storage in the vector database for RAG functionality.
"""

import logging
import hashlib
import tempfile
import os
import time
from typing import List, Dict, Any, Optional, Tuple
from datetime import datetime, timezone
from pathlib import Path

import PyPDF2
from docx import Document as DocxDocument
from pptx import Presentation
from fastapi import UploadFile, HTTPException, status
from langchain_core.documents import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter

from app.schemas.document import (
    DocumentType, DocumentCategory, DocumentUploadResponse,
    DocumentMetadata, KnowledgeBaseStats
)
from app.services.ai.vector_store import get_vector_store_manager
from app.core.database import get_database

logger = logging.getLogger(__name__)


class DocumentService:
    """Service for document processing and knowledge base management"""
    
    def __init__(self):
        self.collection_name = "documents"
        self.max_file_size = 50 * 1024 * 1024  # 50MB
        self.supported_types = {
            "application/pdf": DocumentType.PDF,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": DocumentType.DOCX,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": DocumentType.PPTX,
            "text/plain": DocumentType.TXT
        }
        
        # Text splitter for chunking documents
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
    
    async def upload_document(
        self,
        file: UploadFile,
        category: DocumentCategory,
        uploaded_by: str
    ) -> DocumentUploadResponse:
        """
        Upload and process a document for the knowledge base.
        
        Args:
            file: Uploaded file object
            category: Document category
            uploaded_by: Username of the admin who uploaded the document
            
        Returns:
            DocumentUploadResponse with processing results
            
        Raises:
            HTTPException: If file validation or processing fails
        """
        start_time = time.time()
        
        try:
            # Validate file
            await self._validate_file(file)
            
            # Determine document type
            doc_type = self._get_document_type(file.content_type, file.filename)
            
            # Generate document ID and metadata
            document_id = self._generate_document_id(file.filename, uploaded_by)
            
            # Save file temporarily for processing
            with tempfile.NamedTemporaryFile(delete=False, suffix=f".{doc_type.value}") as temp_file:
                content = await file.read()
                temp_file.write(content)
                temp_file_path = temp_file.name
            
            try:
                # Extract text from document
                text_content, pages_processed = await self._extract_text(temp_file_path, doc_type)
                
                # Create text chunks
                chunks = self._create_chunks(text_content, document_id, file.filename, category)
                
                # Store in vector database
                vectors_stored = await self._store_in_vector_db(chunks)
                
                # Save document metadata
                metadata = DocumentMetadata(
                    document_id=document_id,
                    filename=file.filename,
                    original_filename=file.filename,
                    file_size=len(content),
                    document_type=doc_type,
                    category=category,
                    uploaded_by=uploaded_by,
                    uploaded_at=datetime.now(timezone.utc),
                    pages_processed=pages_processed,
                    chunks_created=len(chunks),
                    vectors_stored=vectors_stored,
                    processing_time=time.time() - start_time,
                    checksum=hashlib.md5(content).hexdigest()
                )
                
                await self._save_metadata(metadata)
                
                logger.info(f"Document processed successfully: {file.filename} -> {vectors_stored} vectors")
                
                return DocumentUploadResponse(
                    document_id=document_id,
                    filename=file.filename,
                    file_size=len(content),
                    document_type=doc_type,
                    category=category,
                    pages_processed=pages_processed,
                    chunks_created=len(chunks),
                    vectors_stored=vectors_stored,
                    processing_time=time.time() - start_time,
                    uploaded_at=metadata.uploaded_at,
                    uploaded_by=uploaded_by
                )
                
            finally:
                # Clean up temporary file
                if os.path.exists(temp_file_path):
                    os.unlink(temp_file_path)
                    
        except Exception as e:
            logger.error(f"Document processing failed: {str(e)}")
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Document processing failed: {str(e)}"
            )
    
    async def _validate_file(self, file: UploadFile) -> None:
        """Validate uploaded file"""
        if not file.filename:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="No filename provided"
            )
        
        if file.content_type not in self.supported_types:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Unsupported file type: {file.content_type}. Supported types: {list(self.supported_types.keys())}"
            )
        
        # Check file size (read a small portion to estimate)
        file.file.seek(0, 2)  # Seek to end
        file_size = file.file.tell()
        file.file.seek(0)  # Reset to beginning
        
        if file_size > self.max_file_size:
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"File too large. Maximum size: {self.max_file_size / (1024*1024):.1f}MB"
            )
    
    def _get_document_type(self, content_type: str, filename: str) -> DocumentType:
        """Determine document type from content type and filename"""
        if content_type in self.supported_types:
            return self.supported_types[content_type]
        
        # Fallback to file extension
        extension = Path(filename).suffix.lower()
        extension_map = {
            ".pdf": DocumentType.PDF,
            ".docx": DocumentType.DOCX,
            ".pptx": DocumentType.PPTX,
            ".txt": DocumentType.TXT
        }
        
        if extension in extension_map:
            return extension_map[extension]
        
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Cannot determine document type for file: {filename}"
        )
    
    def _generate_document_id(self, filename: str, uploaded_by: str) -> str:
        """Generate unique document ID"""
        timestamp = datetime.now(timezone.utc).isoformat()
        content = f"{filename}_{uploaded_by}_{timestamp}"
        return hashlib.sha256(content.encode()).hexdigest()[:16]
    
    async def _extract_text(self, file_path: str, doc_type: DocumentType) -> Tuple[str, Optional[int]]:
        """Extract text content from document"""
        try:
            if doc_type == DocumentType.PDF:
                return await self._extract_pdf_text(file_path)
            elif doc_type == DocumentType.DOCX:
                return await self._extract_docx_text(file_path)
            elif doc_type == DocumentType.PPTX:
                return await self._extract_pptx_text(file_path)
            elif doc_type == DocumentType.TXT:
                return await self._extract_txt_text(file_path)
            else:
                raise ValueError(f"Unsupported document type: {doc_type}")
                
        except Exception as e:
            logger.error(f"Text extraction failed for {file_path}: {str(e)}")
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"Failed to extract text from document: {str(e)}"
            )
    
    async def _extract_pdf_text(self, file_path: str) -> Tuple[str, int]:
        """Extract text from PDF file"""
        text_content = []
        pages_processed = 0
        
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            pages_processed = len(pdf_reader.pages)
            
            for page in pdf_reader.pages:
                text_content.append(page.extract_text())
        
        return "\n\n".join(text_content), pages_processed
    
    async def _extract_docx_text(self, file_path: str) -> Tuple[str, None]:
        """Extract text from DOCX file"""
        doc = DocxDocument(file_path)
        text_content = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text)
        
        return "\n\n".join(text_content), None
    
    async def _extract_pptx_text(self, file_path: str) -> Tuple[str, int]:
        """Extract text from PPTX file"""
        prs = Presentation(file_path)
        text_content = []
        slides_processed = len(prs.slides)
        
        for slide in prs.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text_content.append("\n".join(slide_text))
        
        return "\n\n".join(text_content), slides_processed
    
    async def _extract_txt_text(self, file_path: str) -> Tuple[str, None]:
        """Extract text from TXT file"""
        with open(file_path, 'r', encoding='utf-8') as file:
            return file.read(), None

    def _create_chunks(
        self,
        text_content: str,
        document_id: str,
        filename: str,
        category: DocumentCategory
    ) -> List[Document]:
        """Create text chunks for vector storage"""
        if not text_content.strip():
            logger.warning(f"No text content found in document: {filename}")
            return []

        # Split text into chunks
        chunks = self.text_splitter.split_text(text_content)

        # Create Document objects with metadata
        documents = []
        for i, chunk in enumerate(chunks):
            if chunk.strip():  # Only include non-empty chunks
                doc = Document(
                    page_content=chunk,
                    metadata={
                        "document_id": document_id,
                        "filename": filename,
                        "category": category.value,
                        "chunk_index": i,
                        "source": f"{filename} (chunk {i+1})",
                        "document_type": "knowledge_base"
                    }
                )
                documents.append(doc)

        logger.debug(f"Created {len(documents)} chunks from {filename}")
        return documents

    async def _store_in_vector_db(self, documents: List[Document]) -> int:
        """Store document chunks in vector database"""
        if not documents:
            return 0

        vector_store = get_vector_store_manager()

        if not vector_store._initialized:
            logger.error("Vector store not initialized")
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="Vector database not available"
            )

        success = vector_store.add_documents(documents)

        if not success:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="Failed to store documents in vector database"
            )

        return len(documents)

    async def _save_metadata(self, metadata: DocumentMetadata) -> None:
        """Save document metadata to database"""
        try:
            db = get_database()
            collection = db[self.collection_name]

            document_dict = metadata.model_dump()
            document_dict["_id"] = metadata.document_id

            await collection.insert_one(document_dict)
            logger.debug(f"Saved metadata for document: {metadata.filename}")

        except Exception as e:
            logger.error(f"Failed to save document metadata: {str(e)}")
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="Failed to save document metadata"
            )

    async def get_knowledge_base_stats(self) -> KnowledgeBaseStats:
        """Get statistics about the knowledge base"""
        try:
            db = get_database()
            collection = db[self.collection_name]

            # Get total documents
            total_documents = await collection.count_documents({})

            # Get documents by category
            category_pipeline = [
                {"$group": {"_id": "$category", "count": {"$sum": 1}}}
            ]
            category_results = await collection.aggregate(category_pipeline).to_list(None)
            documents_by_category = {item["_id"]: item["count"] for item in category_results}

            # Get documents by type
            type_pipeline = [
                {"$group": {"_id": "$document_type", "count": {"$sum": 1}}}
            ]
            type_results = await collection.aggregate(type_pipeline).to_list(None)
            documents_by_type = {item["_id"]: item["count"] for item in type_results}

            # Get total size
            size_pipeline = [
                {"$group": {"_id": None, "total_size": {"$sum": "$file_size"}}}
            ]
            size_results = await collection.aggregate(size_pipeline).to_list(None)
            total_size_bytes = size_results[0]["total_size"] if size_results else 0
            total_size_mb = total_size_bytes / (1024 * 1024)

            # Get total vectors (sum of vectors_stored)
            vector_pipeline = [
                {"$group": {"_id": None, "total_vectors": {"$sum": "$vectors_stored"}}}
            ]
            vector_results = await collection.aggregate(vector_pipeline).to_list(None)
            total_vectors = vector_results[0]["total_vectors"] if vector_results else 0

            # Get last updated
            last_doc = await collection.find_one({}, sort=[("uploaded_at", -1)])
            last_updated = last_doc["uploaded_at"] if last_doc else None

            return KnowledgeBaseStats(
                total_documents=total_documents,
                total_vectors=total_vectors,
                documents_by_category=documents_by_category,
                documents_by_type=documents_by_type,
                total_size_mb=round(total_size_mb, 2),
                last_updated=last_updated
            )

        except Exception as e:
            logger.error(f"Failed to get knowledge base stats: {str(e)}")
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail="Failed to retrieve knowledge base statistics"
            )


# Global service instance
document_service = DocumentService()
